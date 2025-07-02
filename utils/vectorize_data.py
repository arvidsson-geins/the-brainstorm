import os
from transformers import AutoTokenizer
from sentence_transformers import SentenceTransformer
import faiss
import pandas as pd
import pickle
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import email
from email import policy

# Initialize embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
MODEL_TOKENIZER_MAPPING = {
    "gpt2": "gpt2",
    "llama3.2": "baseten/Meta-Llama-3-tokenizer",  # Example for LLaMA 2
    "gemma2": "gpt2",      # Replace with actual tokenizer path or name
    "codegemma": "gpt2",
    "wizardlm2": "gpt2",
    "codellama": "gpt2",
}


def regenerate_index_if_missing(agent_name, base_path="./index/data/"):
    """
    Regenerate the FAISS index and data file for the specified agent if missing.

    Args:
        agent_name (str): The name of the agent.
        base_path (str): The base path for index and data files.
    """
    index_file = os.path.join(base_path, f"{agent_name.lower()}_index.bin")
    data_file = os.path.join(base_path, f"{agent_name.lower()}_data.pkl")
    agent_dir = f"./data/{agent_name.lower()}/"
    
    print(f"Regenerating index for {agent_name}...")

    os.makedirs(base_path, exist_ok=True)

    # Ensure the agent directory exists
    if not os.path.exists(agent_dir):
        print(f"Directory for {agent_name} does not exist: {agent_dir}")
        return

    # Load raw data
    raw_data = load_data(agent_name)
    if not raw_data:
        print(f"No data found for {agent_name} in {agent_dir}. Skipping index regeneration.")
        return

    try:
        # Create embeddings and FAISS index
        embeddings = vectorize_data(raw_data)
        index = create_faiss_index(raw_data, embeddings)

        # Save index and data
        save_index(index, index_file)
        save_data(raw_data, data_file)

        print(f"Index and data for {agent_name} saved successfully!")
    except Exception as e:
        print(f"Error regenerating index for {agent_name}: {e}")


def load_data(agent_name):
    """
    Load data from all files in the directory named after the agent.
    
    Args:
        agent_name (str): Name of the agent (e.g., 'Eric').
    
    Returns:
        list: Combined content from all supported files.
    """
    agent_name = agent_name.lower()
    # Directory path for the agent
    agent_dir = f"./data/{agent_name}/"
    combined_data = []

    if not os.path.exists(agent_dir):
        print(f"Directory for {agent_name} does not exist: {agent_dir}")
        return combined_data

    # Iterate through all files in the agent's directory
    for file_name in os.listdir(agent_dir):
        file_path = os.path.join(agent_dir, file_name)
        try:
            # Process CSV files
            if file_name.endswith(".csv"):
                print(f"Processing CSV file: {file_path}")
                csv_data = pd.read_csv(file_path)
                csv_text_data = csv_data.astype(str).agg(' '.join, axis=1).tolist()
                combined_data.extend(csv_text_data)

            # Process Excel files
            elif file_name.endswith((".xls", ".xlsx")):
                print(f"Processing Excel file: {file_path}")
                excel_data = pd.read_excel(file_path)
                excel_text_data = excel_data.astype(str).agg(' '.join, axis=1).tolist()
                combined_data.extend(excel_text_data)

            # Process Markdown files
            elif file_name.endswith(".md"):
                print(f"Processing Markdown file: {file_path}")
                with open(file_path, "r") as f:
                    markdown_data = f.readlines()
                    combined_data.extend(markdown_data)

            # Process plain text files
            elif file_name.endswith(".txt"):
                print(f"Processing Text file: {file_path}")
                with open(file_path, "r") as f:
                    text_file_data = f.readlines()
                    combined_data.extend(text_file_data)

            # Process PDF files
            elif file_name.endswith(".pdf"):
                print(f"Processing PDF file: {file_path}")
                reader = PdfReader(file_path)
                pdf_text_data = [page.extract_text() for page in reader.pages]
                combined_data.extend(pdf_text_data)

            # Process Word documents
            elif file_name.endswith(".docx"):
                print(f"Processing Word file: {file_path}")
                doc = Document(file_path)
                word_text_data = [p.text for p in doc.paragraphs if p.text.strip()]
                combined_data.extend(word_text_data)

            # Process PowerPoint files
            elif file_name.endswith(".pptx"):
                print(f"Processing PowerPoint file: {file_path}")
                presentation = Presentation(file_path)
                ppt_text_data = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            ppt_text_data.append(shape.text)
                combined_data.extend(ppt_text_data)

            # Process Keynote files (if exported to .pptx)
            elif file_name.endswith(".key"):
                print(f"Processing Keynote file: {file_path}")
                # Example: Assume Keynote files are exported to .pptx
                # Convert .key to .pptx manually and process as .pptx

            # Process email files
            elif file_name.endswith(".eml"):
                print(f"Processing Email file: {file_path}")
                with open(file_path, "r") as f:
                    msg = email.message_from_file(f, policy=policy.default)
                    email_text_data = msg.get_body(preferencelist=("plain", "html")).get_content()
                    combined_data.append(email_text_data)

            # Add support for other file types as needed
            else:
                print(f"Skipping unsupported file type: {file_path}")

        except Exception as e:
            print(f"Error processing file {file_path}: {e}")

    return combined_data

# Vectorize the data using the embedding model
def vectorize_data(data):
    embeddings = embedding_model.encode(data)
    return embeddings

# Create a FAISS index from the embeddings
def create_faiss_index(data, embeddings):
    # Initialize FAISS index
    index = faiss.IndexFlatL2(embeddings.shape[1])  # L2 distance metric
    # Add embeddings to the index
    index.add(embeddings)
    return index

# Save the FAISS index to a file
def save_index(index, filename="faiss_index.bin"):
    # Set path to /index/ + filename
    file_path = filename
    faiss.write_index(index, file_path)
    #print(f"Index saved to {file_path}")

# Load the FAISS index from a file
def load_index(filename="faiss_index.bin"):
    # Set path to /index/ + filename
    file_path = filename
    index = faiss.read_index(file_path)
    #print(f"Index loaded from {file_path}")
    return index

# Save the data (e.g., rows of text) for mapping query results
def save_data(data, filename="data.pkl"):
    # Set path to /index/ + filename
    file_path = filename
    with open(file_path, "wb") as f:
        pickle.dump(data, f)
    #print(f"Data saved to {file_path}")

# Load the data for mapping query results
def load_data_file(filename="data.pkl"):
    # Set path to /index/ + filename
    file_path = filename
    with open(file_path, "rb") as f:
        data = pickle.load(f)
    #print(f"Data loaded from {file_path}")
    return data

# Query the FAISS index
def query_index(query, data, index, k=1):
    # Convert the query into an embedding
    query_vector = embedding_model.encode([query])
    # Search for the top-k most similar vectors
    distances, indices = index.search(query_vector, k)
    # Retrieve the matching data rows
    results = [data[i] for i in indices[0]]
    return results

def tokenize_string(content, model_name):
    """
    Tokenize the input string using the tokenizer appropriate for the specified model.

    Args:
        content (str): The input string to tokenize.
        model_name (str): The name of the model to determine the tokenizer.

    Returns:
        list: Tokenized content as a list of integers.
    """
    if model_name not in MODEL_TOKENIZER_MAPPING:
        raise ValueError(f"Tokenizer for model '{model_name}' not found in the mapping.")
    
    tokenizer_path = MODEL_TOKENIZER_MAPPING[model_name]
    tokenizer = AutoTokenizer.from_pretrained(tokenizer_path)
    
    #text = truncate_tokens(content, tokenizer)
    
    text = tokenizer.encode(content, add_special_tokens=False)
    
    # see if text is too long
    #if len(text) > 1024:
    #    text = text[-1024:]
        
    
    return text

def truncate_tokens(content, tokenizer, max_tokens=1024):
    """
    Truncate content to ensure it fits within the token limit.

    Args:
        content (str): Input string to truncate.
        tokenizer: The tokenizer used to compute token length.
        max_tokens (int): Maximum allowed tokens.

    Returns:
        str: Truncated content.
    """
    tokens = tokenizer.encode(content, add_special_tokens=False)
    if len(tokens) > max_tokens:
        truncated_tokens = tokens[-max_tokens:]  # Keep only the last `max_tokens`
        return tokenizer.decode(truncated_tokens)
    return content
